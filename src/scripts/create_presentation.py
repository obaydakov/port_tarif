"""Generate a pitch presentation for the Port Tariff Calculator."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------------------
# Theme colours
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
TEAL = RGBColor(0x00, 0x96, 0x88)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE = RGBColor(0x1E, 0x88, 0xE5)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_ORANGE = RGBColor(0xEF, 0x6C, 0x00)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, size=18, bold=False,
             color=DARK_GRAY, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=16,
                    color=DARK_GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


# =========================================================================
# SLIDE 1 — Title
# =========================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide1, NAVY)

# Accent bar
add_box(slide1, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06), TEAL)

add_text(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "Port Tariff Calculator", size=44, bold=True, color=WHITE)

add_text(slide1, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
         "Agentic RAG System for Maritime Port Dues Automation",
         size=24, color=TEAL)

add_text(slide1, Inches(1), Inches(3.8), Inches(11), Inches(2.5),
         "An AI-powered system that ingests any Port Tariff PDF and a natural language "
         "vessel query to automatically calculate all applicable port dues -- "
         "with no hardcoded rates in the calculation engine.",
         size=18, color=WHITE)

add_text(slide1, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
         "Built with Python  |  OpenAI gpt-5.4  |  LangChain  |  ChromaDB  |  FastAPI",
         size=16, color=RGBColor(0x90, 0xCA, 0xF9))

add_text(slide1, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         "Marcura -- Generative AI Solutions Developer Assessment",
         size=14, color=RGBColor(0x78, 0x78, 0x78))


# =========================================================================
# SLIDE 2 — Architecture & Design Principles
# =========================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, WHITE)

# Header bar
add_box(slide2, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_text(slide2, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
         "Architecture & Design Principles", size=32, bold=True, color=WHITE)

# Left column — pipeline
add_text(slide2, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
         "End-to-End Pipeline", size=20, bold=True, color=NAVY)

pipeline_steps = [
    "1.  Natural Language Query --> LLM parses to structured VesselProfile",
    "2.  PDF --> pdfplumber --> LangChain chunks --> ChromaDB vector store",
    "3.  Semantic retrieval per tariff type (6 parallel queries)",
    "4.  OpenAI structured output --> ExtractedTariffRule (Pydantic)",
    "5.  Generic apply_rule() calculator -- zero tariff-specific code",
    "6.  CalculationResult --> CLI report or JSON API response",
]
add_bullet_list(slide2, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
                pipeline_steps, size=15, color=DARK_GRAY, spacing=Pt(10))

# Vertical divider
add_box(slide2, Inches(6.7), Inches(1.4), Inches(0.03), Inches(5.5), TEAL)

# Right column — principles
add_text(slide2, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.5),
         "Core Principles", size=20, bold=True, color=NAVY)

principles = [
    "Schema-Driven -- LLM populates a typed Pydantic schema; "
    "a single generic engine applies any rule to any vessel",

    "Natural Language Input -- users describe vessels in plain "
    "English; no JSON required",

    "Parallel Extraction -- all 6 tariff types extracted "
    "concurrently via ThreadPoolExecutor (~5x speedup)",

    "Graceful Fallback -- LOW-confidence extractions "
    "transparently use hardcoded rates (same code path)",

    "Rule Caching -- extracted rules persisted by PDF hash; "
    "re-runs skip LLM entirely",

    "Document-Agnostic -- new port tariff = re-index PDF, "
    "zero code changes needed",
]
add_bullet_list(slide2, Inches(7.2), Inches(2.0), Inches(5.5), Inches(5.0),
                principles, size=14, color=DARK_GRAY, spacing=Pt(10))


# =========================================================================
# SLIDE 3 — Solution Design: Modules & Execution Flow
# =========================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, WHITE)

add_box(slide3, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_text(slide3, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
         "Solution Design -- Modules & Execution Flow", size=32, bold=True, color=WHITE)

# --- Left side: Module map ---
add_text(slide3, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
         "Module Responsibilities", size=20, bold=True, color=NAVY)

modules = [
    ("query_parser.py", "NL text -> VesselProfile", ACCENT_BLUE),
    ("document_processor.py", "PDF -> PageContent -> LangChain chunks", ACCENT_BLUE),
    ("tariff_extractor.py", "ChromaDB index + OpenAI structured extraction + disk cache", ACCENT_BLUE),
    ("tariff_schema.py", "ExtractedTariffRule + RateBracket (schema contract)", TEAL),
    ("vessel_profile.py", "VesselProfile Pydantic model + derived props", TEAL),
    ("calculator.py", "Generic apply_rule() -- zero tariff-specific logic", ACCENT_GREEN),
    ("agent.py", "Parallel orchestration + fallback + report formatter", ACCENT_GREEN),
    ("fallback_rates.py", "Transnet hardcoded rates as ExtractedTariffRule objects", ACCENT_ORANGE),
    ("main.py", "CLI: validate / calculate / query", RGBColor(0x78, 0x78, 0x78)),
    ("api.py", "FastAPI: /query, /calculate, /calculate/upload, /health", RGBColor(0x78, 0x78, 0x78)),
]

y = Inches(2.0)
for mod_name, mod_desc, mod_color in modules:
    # Module name box
    box = add_box(slide3, Inches(0.8), y, Inches(2.2), Inches(0.35), mod_color)
    # Name text on the box
    add_text(slide3, Inches(0.85), y, Inches(2.1), Inches(0.35),
             mod_name, size=11, bold=True, color=WHITE)
    # Description beside it
    add_text(slide3, Inches(3.15), y, Inches(3.5), Inches(0.35),
             mod_desc, size=11, color=DARK_GRAY)
    y += Inches(0.42)

# Legend
legend_y = y + Inches(0.2)
legend_items = [
    (ACCENT_BLUE, "AI / LLM layer"),
    (TEAL, "Data models"),
    (ACCENT_GREEN, "Calculation engine"),
    (ACCENT_ORANGE, "Fallback data"),
    (RGBColor(0x78, 0x78, 0x78), "Entry points"),
]
lx = Inches(0.8)
for lc, lt in legend_items:
    add_box(slide3, lx, legend_y, Inches(0.25), Inches(0.2), lc)
    add_text(slide3, lx + Inches(0.3), legend_y, Inches(1.5), Inches(0.2),
             lt, size=9, color=DARK_GRAY)
    lx += Inches(2.0)

# --- Right side: Execution flow diagram (text-based) ---
add_box(slide3, Inches(6.7), Inches(1.4), Inches(0.03), Inches(5.5), TEAL)

add_text(slide3, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.5),
         "Execution Flow", size=20, bold=True, color=NAVY)

# Flow boxes
flow_steps = [
    ("INPUT", "NL query OR structured JSON + Port Tariff PDF", ACCENT_BLUE),
    ("PARSE", "query_parser: LLM structured output -> VesselProfile", ACCENT_BLUE),
    ("INDEX", "document_processor: PDF -> chunks -> ChromaDB", ACCENT_BLUE),
    ("EXTRACT", "tariff_extractor: 6x parallel retrieval + LLM extraction\n"
                "Cache hit? -> skip LLM  |  LOW confidence? -> fallback", ACCENT_GREEN),
    ("CALCULATE", "calculator: apply_rule(vessel, rule) for each tariff\n"
                  "Schema-driven: brackets / per-GT / flat fee / time-based", ACCENT_GREEN),
    ("OUTPUT", "CalculationResult -> CLI report or JSON API response", NAVY),
]

fy = Inches(2.1)
for step_label, step_desc, step_color in flow_steps:
    # Label badge
    add_box(slide3, Inches(7.2), fy, Inches(1.2), Inches(0.35), step_color)
    add_text(slide3, Inches(7.25), fy, Inches(1.15), Inches(0.35),
             step_label, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Description
    txf = add_text(slide3, Inches(8.5), fy, Inches(4.2), Inches(0.7),
                   step_desc, size=11, color=DARK_GRAY)
    fy += Inches(0.82)

    # Arrow between steps (except after last)
    if step_label != "OUTPUT":
        arrow_y = fy - Inches(0.12)
        add_text(slide3, Inches(7.65), arrow_y, Inches(0.5), Inches(0.2),
                 "|", size=14, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
        fy += Inches(0.08)


# =========================================================================
# SLIDE 4 — Accuracy & Validation
# =========================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, WHITE)

add_box(slide4, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_text(slide4, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
         "Accuracy & Validation Results", size=32, bold=True, color=WHITE)

add_text(slide4, Inches(0.8), Inches(1.3), Inches(11), Inches(0.5),
         "Reference vessel: SUDESTADA (Bulk Carrier, GT 51,300) at Durban",
         size=16, color=RGBColor(0x66, 0x66, 0x66))

# Table
rows, cols = 8, 4
tbl_shape = slide4.shapes.add_table(rows, cols,
                                     Inches(1.5), Inches(2.0),
                                     Inches(10), Inches(4.0))
tbl = tbl_shape.table

# Column widths
tbl.columns[0].width = Inches(3.0)
tbl.columns[1].width = Inches(2.5)
tbl.columns[2].width = Inches(2.5)
tbl.columns[3].width = Inches(2.0)

headers = ["Tariff", "Calculated (ZAR)", "Ground Truth (ZAR)", "Error"]
data = [
    ["Light Dues",    "60,062.04",  "60,062.04",  "0.00%"],
    ["VTS Dues",      "33,345.00",  "33,315.75",  "+0.09%"],
    ["Pilotage Dues", "47,189.94",  "47,189.94",  "0.00%"],
    ["Towage Dues",   "147,074.38", "147,074.38", "0.00%"],
    ["Running Lines", "19,854.72",  "19,639.50",  "+1.10%"],
    ["Port Dues",     "199,371.35", "199,549.22", "-0.09%"],
    ["TOTAL",         "506,897.43", "506,830.83", "+0.01%"],
]

def style_cell(cell, text, size=14, bold=False, color=DARK_GRAY, bg=None, align=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Calibri"
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg

for i, h in enumerate(headers):
    style_cell(tbl.cell(0, i), h, size=14, bold=True, color=WHITE, bg=NAVY)

for r, row_data in enumerate(data):
    is_total = r == len(data) - 1
    bg_color = RGBColor(0xE8, 0xF5, 0xE9) if is_total else (LIGHT_GRAY if r % 2 == 0 else WHITE)
    for c, val in enumerate(row_data):
        is_exact = c == 3 and val == "0.00%"
        cell_color = ACCENT_GREEN if is_exact else (NAVY if is_total else DARK_GRAY)
        style_cell(tbl.cell(r + 1, c), val, size=14,
                   bold=is_total or is_exact,
                   color=cell_color, bg=bg_color,
                   align=PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER)

add_text(slide4, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
         "3 tariffs exact match  |  All 6 within +/-1.1%  |  Total within +0.01% of ground truth (spec says 'Approx.')",
         size=15, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


# =========================================================================
# SLIDE 5 — Technical Highlights & Code Quality
# =========================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, WHITE)

add_box(slide5, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_text(slide5, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
         "Technical Highlights & Code Quality", size=32, bold=True, color=WHITE)

# Three columns
col_w = Inches(3.6)
col_gap = Inches(0.4)
col_starts = [Inches(0.8), Inches(0.8) + col_w + col_gap, Inches(0.8) + 2 * (col_w + col_gap)]

# Column 1 — AI / LLM
add_box(slide5, col_starts[0], Inches(1.4), col_w, Inches(0.5), ACCENT_BLUE)
add_text(slide5, col_starts[0], Inches(1.42), col_w, Inches(0.5),
         "  AI & LLM Integration", size=17, bold=True, color=WHITE)

ai_items = [
    "OpenAI structured output -- guaranteed valid Pydantic schema, no JSON parsing",
    "Natural language query parsing -- free-text to VesselProfile via LLM",
    "RAG with ChromaDB -- semantic retrieval with section-aware chunking",
    "Section mismatch detection -- validates extracted sections, auto-downgrades on mismatch",
    "Prompt engineering -- tariff-specific patterns (flat fee, per-GT, bracket) guide the LLM",
]
add_bullet_list(slide5, col_starts[0], Inches(2.1), col_w, Inches(4.5),
                ai_items, size=13, color=DARK_GRAY, spacing=Pt(8))

# Column 2 — Engineering
add_box(slide5, col_starts[1], Inches(1.4), col_w, Inches(0.5), ACCENT_GREEN)
add_text(slide5, col_starts[1], Inches(1.42), col_w, Inches(0.5),
         "  Engineering Quality", size=17, bold=True, color=WHITE)

eng_items = [
    "Strict type safety -- mypy strict + pyright strict + beartype runtime checks",
    "26 unit tests -- all 6 tariff patterns, edge cases, schema validation, error paths",
    "Pydantic v2 models -- typed schemas, model_validator, model_copy",
    "Clean error handling -- failed tariffs produce structured error items, not silent drops",
    "Zero dead code -- no unused deps, no leftover imports",
]
add_bullet_list(slide5, col_starts[1], Inches(2.1), col_w, Inches(4.5),
                eng_items, size=13, color=DARK_GRAY, spacing=Pt(8))

# Column 3 — Performance
add_box(slide5, col_starts[2], Inches(1.4), col_w, Inches(0.5), ACCENT_ORANGE)
add_text(slide5, col_starts[2], Inches(1.42), col_w, Inches(0.5),
         "  Performance & API", size=17, bold=True, color=WHITE)

perf_items = [
    "Parallel extraction -- ThreadPoolExecutor, 6 tariffs concurrently (~5x speedup)",
    "Disk rule cache -- SHA-256 PDF hash keying, skip LLM on cache hit",
    "FastAPI REST -- /query (NL), /calculate (JSON), /calculate/upload (file), /health",
    "Vector store caching -- ChromaDB persisted, rebuilt only on demand",
    "Fallback mode -- instant results with --no-llm, no API key needed",
]
add_bullet_list(slide5, col_starts[2], Inches(2.1), col_w, Inches(4.5),
                perf_items, size=13, color=DARK_GRAY, spacing=Pt(8))


# =========================================================================
# SLIDE 6 — Production Readiness & Roadmap
# =========================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, WHITE)

add_box(slide6, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
add_text(slide6, Inches(0.8), Inches(0.2), Inches(11), Inches(0.7),
         "Production Readiness & Roadmap", size=32, bold=True, color=WHITE)

# Left — What's done
add_text(slide6, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
         "Delivered", size=20, bold=True, color=ACCENT_GREEN)

done_items = [
    "Full agentic RAG pipeline (PDF -> LLM -> calculation)",
    "Natural language vessel query (CLI + API)",
    "Generic schema-driven calculator (no hardcoded logic)",
    "6 tariff types with +0.01% accuracy vs ground truth",
    "26 unit tests, strict type checking (mypy + pyright)",
    "FastAPI with 4 endpoints + Swagger docs",
    "Parallel extraction + disk caching for performance",
    "Graceful fallback for deterministic offline mode",
]
add_bullet_list(slide6, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5),
                done_items, size=14, color=DARK_GRAY, spacing=Pt(8))

# Divider
add_box(slide6, Inches(6.7), Inches(1.4), Inches(0.03), Inches(5.5), TEAL)

# Right — Roadmap
add_text(slide6, Inches(7.2), Inches(1.4), Inches(5.5), Inches(0.5),
         "Roadmap to Production", size=20, bold=True, color=ACCENT_ORANGE)

roadmap_items = [
    "CI/CD -- GitHub Actions (pytest + type-check gates)",
    "Docker -- single-stage build, Cloud Run / ECS / Fargate",
    "Monitoring -- structured logging + OpenTelemetry tracing",
    "Surcharge engine -- time-of-day + event-based surcharges",
    "Multi-document -- index multiple port tariff PDFs",
    "Cargo dues -- Section 7 (dry bulk, breakbulk, containers)",
    "Vessel enrichment -- auto-fetch specs from IMO databases",
    "Multi-currency -- live FX conversion (ZAR/USD/EUR/GBP)",
]
add_bullet_list(slide6, Inches(7.2), Inches(2.0), Inches(5.5), Inches(4.5),
                roadmap_items, size=14, color=DARK_GRAY, spacing=Pt(8))

# Bottom banner
add_box(slide6, Inches(0), Inches(6.6), SLIDE_W, Inches(0.9), NAVY)
add_text(slide6, Inches(0.8), Inches(6.7), Inches(11.5), Inches(0.6),
         "Generalizable by design -- new port tariff = re-index PDF, zero code changes",
         size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)


# =========================================================================
# Save
# =========================================================================
output = "Port_Tariff_Calculator_Pitch.pptx"
prs.save(output)
print(f"Presentation saved: {output}")
